from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "PAU_Market_Tez_Posteri.pptx"

# Poster template requirements from the provided sample:
# - A1 portrait: 594 x 841 mm
# - University logo/name on upper-left
# - Project code on upper-right, white background, Arial 50, black
# - Blue title band, gray two-column content area, GBYF footer strip

# PAU Market brand palette: campus-safe, modern marketplace blue with a
# verification green accent. The template structure stays intact; these colors
# make the poster feel like the actual product rather than a generic sample.
PAU_NAVY = RGBColor(15, 23, 42)
PAU_BLUE = RGBColor(37, 99, 235)
PAU_BLUE_DARK = RGBColor(30, 64, 175)
PAU_SKY = RGBColor(229, 241, 255)
PAU_CLOUD = RGBColor(246, 249, 255)
PAU_GREEN = RGBColor(16, 185, 129)

DARK_BLUE = PAU_NAVY
HEADER_BLUE = PAU_BLUE
SECTION_BLUE = PAU_BLUE_DARK
PAGE_WHITE = RGBColor(255, 255, 255)
PANEL_GRAY = PAU_SKY
PANEL_LIGHT = PAU_CLOUD
BORDER_BLUE = RGBColor(191, 219, 254)
BLACK = RGBColor(0, 0, 0)
TEXT = PAU_NAVY
RED = RGBColor(220, 38, 38)


def cm(value: float):
    return Cm(value)


def inch(value: float):
    return Inches(value)


def set_font(run, size=18, bold=False, color=TEXT, name="Arial"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_rect(slide, x, y, w, h, fill=PAGE_WHITE, line=BLACK, width=0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cm(x), cm(y), cm(w), cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(width)
    return shape


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(cm(x), cm(y), cm(w), cm(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    set_font(p.runs[0], size=size, bold=bold, color=color)
    return box


def add_pill(slide, text, x, y, w, h, fill=PAGE_WHITE, line=BORDER_BLUE, color=TEXT, size=12.5):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cm(x), cm(y), cm(w), cm(h))
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    pill.line.color.rgb = line
    pill.line.width = Pt(0.8)
    pill.text_frame.clear()
    pill.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = pill.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    set_font(p.runs[0], size=size, bold=True, color=color)
    return pill


def add_centered_shape_text(shape, text, size=18, bold=True, color=TEXT):
    shape.text_frame.clear()
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = shape.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    set_font(p.runs[0], size=size, bold=bold, color=color)


def add_paragraph_box(slide, lines, x, y, w, h, size=17, bullet=False):
    box = slide.shapes.add_textbox(cm(x), cm(y), cm(w), cm(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = cm(0.14)
    tf.margin_right = cm(0.12)
    tf.margin_top = cm(0.05)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line if not bullet else f"• {line}"
        p.space_after = Pt(5)
        p.line_spacing = 1.02
        set_font(p.runs[0], size=size, bold=False, color=TEXT)
    return box


def add_section(slide, title, lines, x, y, w, h, body_size=17, bullet=False):
    add_rect(slide, x, y, w, h, fill=PANEL_LIGHT, line=BORDER_BLUE, width=0.6)
    add_rect(slide, x, y, 0.22, h, fill=PAU_GREEN, line=PAU_GREEN, width=0)
    add_rect(slide, x, y, w, 1.18, fill=SECTION_BLUE, line=SECTION_BLUE, width=0.4)
    add_text(slide, title, x + 0.15, y + 0.18, w - 0.3, 0.7, size=22, bold=True, color=PAGE_WHITE, align=PP_ALIGN.CENTER)
    add_paragraph_box(slide, lines, x + 0.33, y + 1.45, w - 0.66, h - 1.65, size=body_size, bullet=bullet)


def add_flow_arrow(slide, x1, y1, x2, y2, color=PAU_BLUE):
    line = slide.shapes.add_connector(1, cm(x1), cm(y1), cm(x2), cm(y2))
    line.line.color.rgb = color
    line.line.width = Pt(1.7)
    line.line.end_arrowhead = True
    return line


def add_architecture_flow(slide, x, y):
    cards = [
        ("React\nArayüz", x, y, 4.8, 2.3, PAU_BLUE),
        (".NET\nAPI", x + 5.65, y, 4.8, 2.3, PAU_BLUE_DARK),
        ("MSSQL\nVeri", x + 11.3, y, 4.8, 2.3, PAU_GREEN),
        ("FastAPI\nÖneri", x + 16.95, y, 5.2, 2.3, PAU_NAVY),
    ]
    for label, cx, cy, cw, ch, fill in cards:
        add_pill(slide, label, cx, cy, cw, ch, fill=fill, line=fill, color=PAGE_WHITE, size=11.5)
    add_flow_arrow(slide, x + 4.8, y + 1.15, x + 5.65, y + 1.15)
    add_flow_arrow(slide, x + 10.45, y + 1.15, x + 11.3, y + 1.15)
    add_flow_arrow(slide, x + 16.1, y + 1.15, x + 16.95, y + 1.15)
    add_text(slide, "Docker Compose ile tek komutla çalışan dağıtık yapı", x, y + 2.6, 22.2, 0.6, size=11.5, bold=True, color=PAU_BLUE_DARK, align=PP_ALIGN.CENTER)


def add_weight_scale(slide, x, y):
    events = [
        ("Görüntüleme", "1.0", PAU_SKY, PAU_BLUE_DARK),
        ("Mesaj", "2.0", RGBColor(219, 234, 254), PAU_BLUE_DARK),
        ("Favori", "3.0", RGBColor(204, 251, 241), PAU_NAVY),
        ("Anlaşma", "4.0", RGBColor(187, 247, 208), PAU_NAVY),
        ("Satıldı", "5.0", PAU_GREEN, PAGE_WHITE),
    ]
    add_text(slide, "Etkileşim gücü", x, y - 0.65, 9.0, 0.5, size=12.5, bold=True, color=PAU_BLUE_DARK)
    step = 3.18
    for idx, (label, weight, fill, color) in enumerate(events):
        px = x + idx * step
        add_pill(slide, weight, px, y, 1.75, 1.12, fill=fill, line=PAU_BLUE, color=color, size=11.5)
        add_text(slide, label, px - 0.7, y + 1.28, 3.15, 0.45, size=8.6, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_flow_arrow(slide, x + 0.8, y + 2.18, x + 14.0, y + 2.18, color=PAU_GREEN)
    add_text(slide, "ilgi zayıf", x, y + 2.55, 3.5, 0.45, size=9.6, bold=False, color=TEXT)
    add_text(slide, "satın alma niyeti güçlü", x + 8.0, y + 2.55, 6.8, 0.45, size=9.6, bold=False, color=TEXT, align=PP_ALIGN.RIGHT)


def add_header_badges(slide, x, y):
    badges = [
        ("PAÜ e-posta doğrulama", PAU_GREEN),
        ("Güvenli mesajlaşma", PAU_NAVY),
        ("Kampüs içi teslim", PAU_BLUE_DARK),
        ("LightFM öneri sistemi", PAU_GREEN),
    ]
    for idx, (label, color) in enumerate(badges):
        add_pill(slide, label, x + idx * 8.4, y, 7.65, 0.85, fill=PAGE_WHITE, line=PAGE_WHITE, color=color, size=8.8)


def add_mini_listing_card(slide, x, y, title, price, category, fill):
    add_rect(slide, x, y, 6.45, 4.55, fill=PAGE_WHITE, line=BORDER_BLUE, width=0.7)
    add_rect(slide, x + 0.25, y + 0.25, 2.05, 2.0, fill=fill, line=fill, width=0.3)
    add_text(slide, "✓", x + 0.68, y + 0.53, 1.25, 0.8, size=21, bold=True, color=PAGE_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + 2.55, y + 0.38, 3.35, 0.75, size=10.7, bold=True, color=PAU_NAVY)
    add_text(slide, price, x + 2.55, y + 1.25, 2.5, 0.65, size=13.5, bold=True, color=PAU_BLUE)
    add_pill(slide, category, x + 2.55, y + 2.15, 2.85, 0.72, fill=PAU_SKY, line=BORDER_BLUE, color=PAU_BLUE_DARK, size=7.6)
    add_text(slide, "PAÜ Kampüsü", x + 2.55, y + 3.1, 2.9, 0.45, size=8.5, bold=False, color=RGBColor(71, 85, 105))
    add_text(slide, "♥", x + 5.35, y + 0.2, 0.65, 0.55, size=13, bold=True, color=RED, align=PP_ALIGN.CENTER)


def add_marketplace_mockup(slide, x, y):
    add_rect(slide, x, y, 21.4, 8.2, fill=PAGE_WHITE, line=BORDER_BLUE, width=0.8)
    add_rect(slide, x, y, 21.4, 1.35, fill=PAU_NAVY, line=PAU_NAVY, width=0.3)
    add_text(slide, "PAÜ Market arayüz örneği", x + 0.7, y + 0.32, 8.5, 0.6, size=12, bold=True, color=PAGE_WHITE)
    add_pill(slide, "Ara: macbook, şarj, ders notu", x + 10.6, y + 0.32, 9.3, 0.72, fill=PAGE_WHITE, line=PAGE_WHITE, color=PAU_NAVY, size=8.6)
    add_mini_listing_card(slide, x + 0.75, y + 2.05, "MacBook Air", "₺18.500", "Elektronik", PAU_BLUE)
    add_mini_listing_card(slide, x + 7.45, y + 2.05, "Şarj Aleti", "₺250", "Elektronik", PAU_GREEN)
    add_mini_listing_card(slide, x + 14.15, y + 2.05, "Algoritma Notu", "₺50", "Not / Özet", PAU_BLUE_DARK)


def add_matrix_heatmap(slide, x, y):
    add_rect(slide, x, y, 8.2, 5.5, fill=PAGE_WHITE, line=BORDER_BLUE, width=0.7)
    add_text(slide, "R[u,i] Etkileşim Matrisi", x + 0.3, y + 0.25, 7.5, 0.5, size=10.7, bold=True, color=PAU_BLUE_DARK, align=PP_ALIGN.CENTER)
    values = [
        [1, 0, 3, 0, 5],
        [0, 2, 0, 4, 0],
        [3, 0, 1, 0, 4],
        [0, 5, 2, 0, 0],
    ]
    colors = {
        0: RGBColor(226, 232, 240),
        1: RGBColor(191, 219, 254),
        2: RGBColor(147, 197, 253),
        3: RGBColor(96, 165, 250),
        4: RGBColor(52, 211, 153),
        5: PAU_GREEN,
    }
    start_x = x + 1.45
    start_y = y + 1.25
    cell = 0.78
    for c in range(5):
        add_text(slide, f"İ{c+1}", start_x + c * 1.02, y + 0.85, 0.75, 0.35, size=7.4, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    for r, row in enumerate(values):
        add_text(slide, f"U{r+1}", x + 0.55, start_y + r * 0.95 + 0.12, 0.55, 0.35, size=7.4, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        for c, value in enumerate(row):
            add_rect(slide, start_x + c * 1.02, start_y + r * 0.95, cell, cell, fill=colors[value], line=PAGE_WHITE, width=0.35)
    add_text(slide, "0: yok  •  5: güçlü niyet", x + 0.65, y + 4.85, 6.9, 0.35, size=7.4, bold=True, color=RGBColor(71, 85, 105), align=PP_ALIGN.CENTER)


def add_pilot_bar_chart(slide, x, y):
    add_rect(slide, x, y, 21.8, 5.2, fill=PAGE_WHITE, line=BORDER_BLUE, width=0.7)
    add_text(slide, "Pilot veri hedefi", x + 0.6, y + 0.35, 7.0, 0.55, size=11.5, bold=True, color=PAU_BLUE_DARK)
    bars = [
        ("İlan", 50, PAU_BLUE),
        ("Öğrenci", 20, PAU_GREEN),
        ("Etkileşim", 400, PAU_NAVY),
    ]
    max_value = 400
    for idx, (label, value, color) in enumerate(bars):
        by = y + 1.25 + idx * 1.18
        add_text(slide, label, x + 0.7, by + 0.12, 3.0, 0.35, size=8.8, bold=True, color=TEXT)
        add_rect(slide, x + 4.0, by, 14.6, 0.55, fill=PAU_SKY, line=PAU_SKY, width=0.2)
        add_rect(slide, x + 4.0, by, max(1.1, 14.6 * value / max_value), 0.55, fill=color, line=color, width=0.2)
        add_text(slide, str(value), x + 19.0, by - 0.02, 1.8, 0.45, size=9.2, bold=True, color=color, align=PP_ALIGN.RIGHT)
    add_text(slide, "Sentetik veri değil, kontrollü gerçek kampüs etkileşimi", x + 0.7, y + 4.35, 20.0, 0.45, size=8.8, bold=True, color=RGBColor(71, 85, 105), align=PP_ALIGN.CENTER)


def add_metric_cards(slide, x, y):
    cards = [
        ("Precision@5", "İlk 5 öneri\nne kadar isabetli?"),
        ("Recall@5", "İlgili ilanların\nne kadarı yakalandı?"),
        ("NDCG@5", "Doğru ilanlar\nüst sırada mı?"),
    ]
    for idx, (title, desc) in enumerate(cards):
        cx = x + idx * 7.35
        add_rect(slide, cx, y, 6.65, 3.3, fill=PAGE_WHITE, line=BORDER_BLUE, width=0.8)
        add_rect(slide, cx, y, 6.65, 0.24, fill=PAU_GREEN if idx == 0 else PAU_BLUE, line=PAU_GREEN if idx == 0 else PAU_BLUE, width=0)
        add_text(slide, title, cx + 0.35, y + 0.55, 5.95, 0.5, size=10.2, bold=True, color=PAU_BLUE_DARK, align=PP_ALIGN.CENTER)
        add_text(slide, desc, cx + 0.55, y + 1.35, 5.55, 1.05, size=8.4, bold=False, color=TEXT, align=PP_ALIGN.CENTER)


def add_table(slide, x, y, w, h):
    rows, cols = 6, 3
    table_shape = slide.shapes.add_table(rows, cols, cm(x), cm(y), cm(w), cm(h))
    table = table_shape.table
    table.columns[0].width = cm(w * 0.36)
    table.columns[1].width = cm(w * 0.25)
    table.columns[2].width = cm(w * 0.39)

    values = [
        ["Metrik", "Durum", "Anlamı"],
        ["Precision@5", "Pilot sonrası", "İlk 5 önerinin isabeti"],
        ["HitRate@5", "Pilot sonrası", "İlk 5'te doğru öneri var mı?"],
        ["Recall@5", "Pilot sonrası", "İlgili ilanları yakalama oranı"],
        ["NDCG@5", "Pilot sonrası", "Doğru önerinin sıralamadaki yeri"],
        ["RMSE", "Pilot sonrası", "Etkileşim tahmin hatası"],
    ]

    for r, row in enumerate(values):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = SECTION_BLUE if r == 0 else (PAU_SKY if r % 2 else PAGE_WHITE)
            p = cell.text_frame.paragraphs[0]
            p.text = text
            set_font(p.runs[0], size=12.5, bold=(r == 0 or c == 0), color=PAGE_WHITE if r == 0 else TEXT)


def add_university_seal(slide, x, y, size_cm=5.0):
    # The sample states the university logo should be at most 5 cm high.
    outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, cm(x), cm(y), cm(size_cm), cm(size_cm))
    outer.fill.solid()
    outer.fill.fore_color.rgb = PAGE_WHITE
    outer.line.color.rgb = DARK_BLUE
    outer.line.width = Pt(1.5)

    inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, cm(x + 0.48), cm(y + 0.48), cm(size_cm - 0.96), cm(size_cm - 0.96))
    inner.fill.solid()
    inner.fill.fore_color.rgb = PAU_SKY
    inner.line.color.rgb = DARK_BLUE
    inner.line.width = Pt(0.8)

    add_text(slide, "PAÜ", x + 0.65, y + 1.45, size_cm - 1.3, 0.9, size=25, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, "DENİZLİ", x + 0.65, y + 2.6, size_cm - 1.3, 0.55, size=13, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)


def add_gbyf_mark(slide, x, y):
    logo = slide.shapes.add_shape(MSO_SHAPE.OVAL, cm(x), cm(y), cm(2.4), cm(2.4))
    logo.fill.solid()
    logo.fill.fore_color.rgb = PAU_BLUE
    logo.line.color.rgb = RGBColor(235, 235, 235)
    logo.line.width = Pt(1.5)
    add_text(slide, "GBYF", x + 0.25, y + 0.78, 1.9, 0.55, size=15, bold=True, color=PAGE_WHITE, align=PP_ALIGN.CENTER)


def build_poster():
    prs = Presentation()
    prs.slide_width = cm(59.4)
    prs.slide_height = cm(84.1)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # White A1 page.
    add_rect(slide, 0, 0, 59.4, 84.1, fill=PAGE_WHITE, line=PAGE_WHITE, width=0)

    page_x = 2.15
    page_y = 1.55
    page_w = 55.1
    page_h = 80.95
    add_rect(slide, page_x, page_y, page_w, page_h, fill=PAGE_WHITE, line=RGBColor(210, 210, 210), width=0.8)

    # Header band.
    header_y = page_y + 0.35
    header_h = 10.2
    add_rect(slide, page_x + 0.15, header_y, page_w - 0.3, header_h, fill=HEADER_BLUE, line=HEADER_BLUE, width=0.6)
    add_university_seal(slide, page_x + 0.7, header_y + 0.55, 5.0)

    # Required top-right project code box: Arial 50, black, white background.
    code_box = add_rect(slide, page_x + page_w - 14.65, header_y, 14.5, 2.6, fill=PAGE_WHITE, line=PAGE_WHITE, width=0.4)
    add_centered_shape_text(code_box, "GBYF Proje No", size=50, bold=True, color=BLACK)

    title_x = page_x + 6.4
    title_w = page_w - 21.5
    add_text(slide, "Pamukkale Üniversitesi", title_x, header_y + 0.75, title_w, 1.1, size=31, bold=True, color=PAGE_WHITE, align=PP_ALIGN.CENTER)
    add_text(
        slide,
        "PAÜ Market: Kampüs İçi Akıllı",
        title_x,
        header_y + 2.25,
        title_w,
        1.25,
        size=32,
        bold=True,
        color=PAGE_WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "İkinci El Pazar Yeri",
        title_x,
        header_y + 3.35,
        title_w,
        1.25,
        size=32,
        bold=True,
        color=PAGE_WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "Hasan Toğmuş, Berke  --  Danışman: [Danışman Öğretim Üyesi]",
        title_x,
        header_y + 6.35,
        title_w,
        0.8,
        size=16,
        bold=True,
        color=PAGE_WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "Bilgisayar Mühendisliği Bitirme Projesi",
        title_x,
        header_y + 7.35,
        title_w,
        0.75,
        size=15,
        bold=True,
        color=RGBColor(235, 242, 250),
        align=PP_ALIGN.CENTER,
    )
    add_header_badges(slide, title_x + 0.2, header_y + 8.55)

    # Main gray content area.
    content_y = header_y + header_h + 0.7
    content_h = 63.2
    add_rect(slide, page_x + 0.95, content_y, page_w - 1.9, content_h, fill=PANEL_GRAY, line=BLACK, width=0.9)

    left_x = page_x + 1.55
    right_x = page_x + 28.45
    col_w = 25.2
    gutter_x = page_x + 27.15
    add_rect(slide, gutter_x, content_y + 0.4, 0.35, content_h - 0.8, fill=PAGE_WHITE, line=PAGE_WHITE, width=0)

    add_section(
        slide,
        "Proje Özeti",
        [
            "PAÜ Market, Pamukkale Üniversitesi öğrencilerine özel kapalı bir C2C ikinci el pazar yeri uygulamasıdır.",
            "Okul e-postasıyla doğrulanan öğrenciler ilan ekler, ilanları keşfeder, favoriler, satıcıya mesaj atar ve anlaşma isteği gönderir.",
            "Tezin teknik odağı, kullanıcı etkileşimlerinden öğrenen kişiselleştirilmiş öneri sistemidir.",
        ],
        left_x,
        content_y + 0.65,
        col_w,
        10.0,
        body_size=16.2,
    )
    add_pill(slide, "Sadece PAÜ öğrencileri", left_x + 0.85, content_y + 8.85, 7.0, 1.0, fill=PAGE_WHITE, line=PAU_GREEN, color=PAU_NAVY, size=10.8)
    add_pill(slide, "Kampüs içi alışveriş", left_x + 8.4, content_y + 8.85, 6.7, 1.0, fill=PAGE_WHITE, line=PAU_BLUE, color=PAU_NAVY, size=10.8)
    add_pill(slide, "Akıllı öneriler", left_x + 15.65, content_y + 8.85, 5.8, 1.0, fill=PAGE_WHITE, line=PAU_BLUE, color=PAU_NAVY, size=10.8)

    add_section(
        slide,
        "Problem ve Amaç",
        [
            "Genel e-ticaret platformları öğrenciler için güven, yerellik ve erişim problemini tam çözmemektedir.",
            "Bu projede amaç; sadece PAÜ öğrencilerine açık, doğrulanmış, moderasyon destekli ve öneri sistemiyle kişiselleşen bir kampüs pazarı geliştirmektir.",
            "Platform güvenli kullanıcı doğrulama, ilan yönetimi, mesajlaşma, anlaşma/satış kaydı ve satıcı değerlendirme süreçlerini tek sistemde toplar.",
        ],
        left_x,
        content_y + 11.55,
        col_w,
        11.0,
        body_size=16.0,
    )

    add_section(
        slide,
        "Sistem Mimarisi ve Teknolojiler",
        [
            "Frontend: React + Vite ile responsive kullanıcı arayüzü.",
            "Backend: ASP.NET Core Web API, JWT kimlik doğrulama, servis katmanı ve DTO yapısı.",
            "Veritabanı: MSSQL + Entity Framework Core.",
            "Öneri Servisi: Python + FastAPI mikroservisi.",
            "DevOps: Docker Compose ile frontend, backend, SQL ve recommender birlikte çalıştırılır.",
        ],
        left_x,
        content_y + 23.45,
        col_w,
        12.4,
        body_size=15.8,
        bullet=True,
    )
    add_architecture_flow(slide, left_x + 1.35, content_y + 31.7)

    add_section(
        slide,
        "Platform Özellikleri",
        [
            "PAÜ e-posta doğrulamalı kayıt, giriş ve şifre sıfırlama",
            "Çoklu fotoğraflı ilan oluşturma ve kapak görseli sıralama",
            "Favori, güvenli mesajlaşma ve anlaşma isteği",
            "Satıldı akışı, alıcı kaydı ve satıcı değerlendirme",
            "Admin paneli ve ilan moderasyonu",
        ],
        left_x,
        content_y + 36.75,
        col_w,
        21.4,
        body_size=16.0,
        bullet=True,
    )
    add_marketplace_mockup(slide, left_x + 1.15, content_y + 47.25)

    add_section(
        slide,
        "Öneri Sistemi",
        [
            "Öneri sistemi, kullanıcı-ilan etkileşimlerini implicit feedback olarak yorumlar.",
            "Etkileşim ağırlıkları: görüntüleme=1.0, mesaj=2.0, favori=3.0, anlaşma isteği=4.0, kabul=4.5, satın alma/satıldı=5.0.",
            "Matris yapısı R[u, i] biçimindedir; satırlar kullanıcıları, sütunlar ilanları, hücreler ilgi gücünü temsil eder.",
            "LightFM tabanlı hibrit yaklaşım collaborative sinyalleri kategori/içerik özellikleriyle birleştirir.",
        ],
        right_x,
        content_y + 0.65,
        col_w,
        13.2,
        body_size=15.9,
    )
    add_weight_scale(slide, right_x + 1.0, content_y + 8.85)
    add_matrix_heatmap(slide, right_x + 16.0, content_y + 4.35)

    add_section(
        slide,
        "Pilot Veri Planı",
        [
            "Tez sunumu öncesinde küçük ölçekli gerçek kampüs pilotu yapılacaktır.",
            "Planlanan veri: 30-50 ilan, 10-20 PAÜ öğrencisi, kişi başı 10-20 etkileşim.",
            "Her kullanıcı kendi okul e-postasıyla giriş yapacak; görüntüleme, favori, mesaj ve anlaşma isteği davranışları toplanacaktır.",
            "Amaç sentetik veri yerine gerçek öğrencilerden kontrollü ve açıklanabilir pilot veri elde etmektir.",
        ],
        right_x,
        content_y + 14.75,
        col_w,
        11.6,
        body_size=15.9,
    )
    add_pilot_bar_chart(slide, right_x + 1.4, content_y + 20.45)

    add_section(
        slide,
        "Değerlendirme Yaklaşımı",
        [
            "Bu sistem klasik sınıflandırma problemi olmadığı için tek başına accuracy metriği kullanılmaz.",
            "Recommender sistemlerde Top-N sıralama metrikleri daha uygundur: Precision@5, Recall@5, HitRate@5, NDCG@5 ve RMSE.",
            "Pilot veri toplandıktan sonra model PAÜ Market verisiyle yeniden eğitilecek ve sonuçlar bu tabloya işlenecektir.",
        ],
        right_x,
        content_y + 27.25,
        col_w,
        10.0,
        body_size=15.4,
    )
    add_table(slide, right_x + 0.5, content_y + 37.75, col_w - 1.0, 7.5)
    add_metric_cards(slide, right_x + 0.75, content_y + 45.75)

    add_section(
        slide,
        "Beklenen Sonuç ve Katkılar",
        [
            "PAÜ Market, kampüs içinde doğrulanmış ve güvenli ikinci el alışveriş ortamı sunar.",
            "Mesajlaşma ve anlaşma akışı sayesinde gerçek satın alma niyeti veri olarak toplanır.",
            "Öneri sistemi kullanıcıların ilgi alanlarına göre ilanları sıralayarak keşif sürecini kolaylaştırır.",
            "Mikroservis mimarisi, öneri modelini ana backend'den ayırarak sürdürülebilir bir yapı sağlar.",
        ],
        right_x,
        content_y + 50.25,
        col_w,
        9.0,
        body_size=15.6,
        bullet=True,
    )

    # Footer: approximately 3 cm high strip, matching the sample poster's bottom identity area.
    footer_y = page_y + page_h - 3.35
    add_gbyf_mark(slide, page_x + 11.0, footer_y + 0.45)
    add_text(slide, "Genç Beyinler Yeni Fikirler", page_x + 14.0, footer_y + 0.35, 24.0, 0.9, size=20, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    add_text(slide, "Proje Pazarı ve Bitirme Projeleri Ortak Sergisi", page_x + 14.0, footer_y + 1.35, 24.0, 0.75, size=16, bold=False, color=BLACK, align=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f"Poster saved: {OUTPUT}")


if __name__ == "__main__":
    build_poster()
